import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def process_shape(shape, f):
    # Check if shape is group
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        f.write(f"\n[Group Shape: {shape.name}]\n")
        for sub_shape in shape.shapes:
            process_shape(sub_shape, f)
        return
        
    if shape.has_text_frame:
        text = shape.text.strip()
        if text:
            f.write(f"[{shape.name} Text]: {text}\n")
            
    if shape.has_table:
        f.write(f"\n[{shape.name} Table]:\n")
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            row_cells = []
            for col_idx, cell in enumerate(row.cells):
                cell_text = " ".join([p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()])
                row_cells.append(cell_text)
            f.write(" | ".join(row_cells) + "\n")
        f.write("\n")

def extract_presentation_details_advanced(pptx_path, output_path):
    prs = Presentation(pptx_path)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(f"Presentation: {os.path.basename(pptx_path)}\n")
        f.write(f"Number of slides: {len(prs.slides)}\n\n")
        
        for i, slide in enumerate(prs.slides, start=1):
            f.write(f"==================================================\n")
            f.write(f"SLIDE {i}\n")
            f.write(f"==================================================\n\n")
            
            f.write("--- SHAPES & TEXTS ---\n")
            for shape in slide.shapes:
                process_shape(shape, f)
            
            # Extract notes
            f.write("\n--- SPEAKER NOTES ---\n")
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    f.write(notes_text + "\n")
                else:
                    f.write("(None)\n")
            else:
                f.write("(None)\n")
            
            f.write("\n\n")

if __name__ == '__main__':
    pptx_path = '/Users/keith/Documents/Antigravity Project/ppt/FNA Training 2026.pptx'
    output_path = '/Users/keith/Documents/Antigravity Project/scratch/pptx_full_details_advanced.txt'
    extract_presentation_details_advanced(pptx_path, output_path)
    print(f"Successfully extracted details to {output_path}")
