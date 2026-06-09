from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def inspect_slides(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, start=1):
        print(f"\nSlide {i}:")
        for j, shape in enumerate(slide.shapes, start=1):
            text_preview = ""
            if shape.has_text_frame and shape.text.strip():
                text_preview = shape.text.strip().replace('\n', ' ')[:60]
                if len(shape.text) > 60:
                    text_preview += "..."
            
            print(f"  Shape {j}: Name='{shape.name}', Type={shape.shape_type}, Text='{text_preview}'")
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for k, sub_shape in enumerate(shape.shapes, start=1):
                    sub_text = ""
                    if sub_shape.has_text_frame and sub_shape.text.strip():
                        sub_text = sub_shape.text.strip().replace('\n', ' ')[:40]
                    print(f"    Sub-Shape {k}: Name='{sub_shape.name}', Type={sub_shape.shape_type}, Text='{sub_text}'")

if __name__ == '__main__':
    pptx_path = '/Users/keith/Documents/Antigravity Project/ppt/FNA Training 2026.pptx'
    inspect_slides(pptx_path)
